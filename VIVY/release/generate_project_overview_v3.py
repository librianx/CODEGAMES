from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = Path(r"d:\lib\CODEGAMES\VIVY\release\VIVY_Project_Overview_v3.pptx")
AVATAR_PATH = ROOT / "static" / "images" / "VIVYstatr.png"

BG = RGBColor(8, 14, 23)
PANEL = RGBColor(18, 28, 42)
PANEL_2 = RGBColor(23, 38, 56)
PANEL_3 = RGBColor(15, 44, 52)
ACCENT = RGBColor(51, 194, 232)
ACCENT_2 = RGBColor(125, 233, 255)
GREEN = RGBColor(91, 206, 151)
WARN = RGBColor(255, 188, 96)
TEXT = RGBColor(239, 248, 255)
TEXT_DIM = RGBColor(176, 197, 218)
MUTED = RGBColor(97, 126, 150)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, x, y, w, h, text, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, main: str, sub: str = ""):
    add_textbox(slide, 0.68, 0.38, 11.6, 0.55, main, size=25, bold=True)
    if sub:
        add_textbox(slide, 0.72, 0.92, 11.6, 0.34, sub, size=11, color=TEXT_DIM)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.72),
        Inches(1.28),
        Inches(11.9),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT


def footer(slide, label: str = "VIVY 桌宠项目汇报"):
    add_textbox(slide, 0.72, 6.98, 11.9, 0.22, label, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def panel(slide, x, y, w, h, fill=PANEL, line=ACCENT_2):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.1)
    return shape


def bullet_panel(slide, heading: str, bullets: list[str], x, y, w, h, size=15, fill=PANEL):
    shape = panel(slide, x, y, w, h, fill=fill)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(15)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(8)

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = heading
    r0.font.name = "Microsoft YaHei"
    r0.font.size = Pt(18)
    r0.font.bold = True
    r0.font.color.rgb = ACCENT_2
    p0.space_after = Pt(8)

    for item in bullets:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)
    return shape


def metric(slide, x, y, w, h, number: str, label: str, color=ACCENT):
    shape = panel(slide, x, y, w, h, fill=PANEL_2, line=color)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_DIM


def flow(slide, steps: list[str], x, y, w, step_h=0.48, gap=0.15, color=ACCENT):
    top = y
    for i, step in enumerate(steps, start=1):
        shape = panel(slide, x, top, w, step_h, fill=PANEL_2 if i % 2 else PANEL, line=color)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i}. {step}"
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = TEXT
        top += step_h + gap
        if i < len(steps):
            arr = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                Inches(x + w / 2 - 0.14),
                Inches(top - gap + 0.02),
                Inches(0.28),
                Inches(gap + 0.06),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = color
            arr.line.color.rgb = color


def swimlane(slide, rows: list[tuple[str, str]], x, y, w):
    for idx, (left, right) in enumerate(rows):
        yy = y + idx * 0.72
        label = panel(slide, x, yy, 2.15, 0.48, fill=PANEL_3, line=ACCENT)
        label.text_frame.clear()
        label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = label.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = left
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = ACCENT_2

        body = panel(slide, x + 2.35, yy, w - 2.35, 0.48, fill=PANEL if idx % 2 else PANEL_2, line=MUTED)
        body.text_frame.clear()
        body.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = body.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = right
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(12)
        r2.font.color.rgb = TEXT


def make_slide(prs, main: str, sub: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, main, sub)
    footer(slide)
    return slide


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.68),
        Inches(0.72),
        Inches(0.12),
        Inches(5.9),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.color.rgb = ACCENT

    add_textbox(slide, 1.03, 0.82, 7.2, 0.66, "VIVY 桌宠项目汇报", size=31, bold=True)
    add_textbox(
        slide,
        1.06,
        1.54,
        8.0,
        0.45,
        "陪伴型 AI 桌宠｜对话记忆｜天气问候｜创作辅助｜双平台交付",
        size=14,
        color=TEXT_DIM,
    )
    bullet_panel(
        slide,
        "项目定位",
        [
            "把 AI 对话、情绪陪伴、创作辅助和本地记忆做成桌面常驻角色。",
            "用户无需打开网页，通过桌宠即可聊天、获取提醒、读文档、进入写作辅助。",
            "强调“像在身边”的低压力交互，而不是一次性工具式问答。",
        ],
        1.05,
        2.45,
        7.2,
        2.35,
        size=16,
    )
    metric(slide, 1.05, 5.15, 1.75, 0.9, "AI", "本地桌宠入口")
    metric(slide, 3.05, 5.15, 1.75, 0.9, "记忆", "长期陪伴状态", color=GREEN)
    metric(slide, 5.05, 5.15, 1.75, 0.9, "双平台", "Windows / macOS", color=WARN)
    metric(slide, 7.05, 5.15, 1.75, 0.9, "创作", "文档与写作辅助")

    if AVATAR_PATH.exists():
        slide.shapes.add_picture(str(AVATAR_PATH), Inches(8.75), Inches(1.0), height=Inches(5.8))
    footer(slide, "VIVY Project Overview · 2026")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)

    slide = make_slide(prs, "1. 作品简介", "VIVY 是一个以桌宠形态运行的本地 AI 交互系统。")
    bullet_panel(
        slide,
        "项目是什么",
        [
            "桌面常驻：透明无边框、可拖拽、可待机收起。",
            "AI 交互：聊天、图片理解入口、语音输入、语音回放。",
            "陪伴表达：天气问候、时间感知、情绪安慰、唱歌互动。",
            "创作辅助：灵感、短剧、读取文档、沉浸写作窗口。",
        ],
        0.8,
        1.65,
        5.65,
        4.8,
    )
    bullet_panel(
        slide,
        "项目目标",
        [
            "让 AI 从“网页里的工具”变成桌面上的长期角色。",
            "让用户更自然地表达需求、情绪和创作想法。",
            "在本地保存偏好与状态，让 Vivy 记得用户和使用节奏。",
            "形成可演示、可打包、可继续扩展的桌宠产品原型。",
        ],
        6.75,
        1.65,
        5.75,
        4.8,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "2. 产品作用", "VIVY 同时承担陪伴、提醒、创作入口和轻量助手四类作用。")
    bullet_panel(
        slide,
        "用户价值",
        [
            "陪伴：降低用户开口成本，让桌面上有一个持续回应的角色。",
            "提醒：根据时间、天气、打开次数和每日时段提供适度问候。",
            "效率：把聊天、文档、写作、记忆和语音集中到一个桌宠入口。",
            "创作：给写作者提供灵感、续写、润色、点评和文档参考能力。",
        ],
        0.8,
        1.55,
        6.1,
        4.95,
    )
    bullet_panel(
        slide,
        "典型场景",
        [
            "早上第一次打开电脑，Vivy 播报天气和防护提醒。",
            "学习或工作间隙，用一句话快速问问题、看时间或查天气。",
            "心情低落时，先被接住情绪，再选择聊天或听歌。",
            "写作时读取参考文档，在沉浸窗口中获得创作辅助。",
        ],
        7.15,
        1.55,
        5.25,
        4.95,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "3. 用户体验设计", "体验重点是轻量、可控、可恢复，以及更像真实陪伴对象。")
    swimlane(
        slide,
        [
            ("轻量入口", "双击 start_all 或 exe 启动，桌宠直接常驻桌面。"),
            ("自然表达", "支持文字、语音、图片和快捷按钮，不必切换多个工具。"),
            ("状态可控", "语音回放、待机时长、记忆面板、位置与天气设置可手动调整。"),
            ("不打扰", "每日首次问候和时段问候都有记忆去重，关闭重开不会重复寒暄。"),
            ("可打断", "新问题可以中断旧语音和旧回复，减少串线。"),
            ("有温度", "唱歌不是情绪关键词触发，而是先安慰、再邀请、再由用户选择。"),
        ],
        0.95,
        1.7,
        11.4,
    )

    slide = make_slide(prs, "4. 功能说明", "目前已形成较完整的桌宠功能闭环。")
    bullet_panel(
        slide,
        "基础能力",
        [
            "普通聊天与创作模式切换。",
            "时间查询：直接读取电脑本地时间，避免模型猜错。",
            "天气查询：手动默认城市 + 可选 IP 粗定位。",
            "图片理解入口：可接入支持视觉的模型。",
            "语音输入和语音回放开关。",
        ],
        0.75,
        1.55,
        3.95,
        5.1,
    )
    bullet_panel(
        slide,
        "陪伴能力",
        [
            "每天第一次打开自动问候。",
            "早 / 中 / 晚 / 深夜时段自动问候。",
            "早上问候会播报天气并给防护建议。",
            "情绪低落先安慰，明确同意后再唱歌。",
            "每日状态写入本地记忆，避免重复打扰。",
        ],
        4.95,
        1.55,
        3.95,
        5.1,
        fill=PANEL_2,
    )
    bullet_panel(
        slide,
        "创作能力",
        [
            "今日灵感和脑洞短剧。",
            "读取 txt / md / docx / pdf 文档。",
            "沉浸写作：新建、保存、全屏、目标字数。",
            "润色、续写、点评、加强和自定义指令。",
            "最近文档和自动备份恢复。",
        ],
        9.15,
        1.55,
        3.4,
        5.1,
    )

    slide = make_slide(prs, "5. 记忆与每日问候", "Vivy 会记住今天已经见过用户几次、哪些时段已经问候过。")
    bullet_panel(
        slide,
        "记忆字段",
        [
            "daily_greeting_state：按日期保存每日问候状态。",
            "open_count：当天打开 Vivy 的次数。",
            "open_greeted：当天首次打开问候是否已完成。",
            "slots：早 / 中 / 晚 / 深夜哪些时段已问候。",
            "open_event_ids：防止同一次启动重试时重复计数。",
        ],
        0.85,
        1.55,
        5.45,
        5.2,
    )
    flow(
        slide,
        [
            "桌宠启动生成本次 open_event_id",
            "调用 /api/weather/greeting trigger=open",
            "后端读取 daily_greeting_state",
            "第一次打开则问候并标记 open_greeted",
            "关闭重开只增加 open_count，不重复问候",
            "定时问候用 trigger=schedule 按 slot 去重",
        ],
        6.85,
        1.55,
        5.25,
        step_h=0.5,
        gap=0.18,
        color=GREEN,
    )

    slide = make_slide(prs, "6. 时间、天气与位置系统", "Vivy 能理解当前时间，也能围绕用户所在地提供天气提醒。")
    bullet_panel(
        slide,
        "时间能力",
        [
            "前端和后端都支持本地时间直答。",
            "回答包含日期、星期、时区和时间段。",
            "避免模型根据训练数据或服务器环境猜错时间。",
        ],
        0.85,
        1.65,
        3.75,
        4.85,
    )
    bullet_panel(
        slide,
        "天气能力",
        [
            "支持“今天会下雨吗”“上海明天几度”等自然提问。",
            "城市解析支持手动默认城市和消息中的显式城市。",
            "可选 IP 粗定位，仅用于天气，不做持续定位。",
            "早晨问候会结合天气生成防护建议。",
        ],
        4.9,
        1.65,
        3.95,
        4.85,
        fill=PANEL_2,
    )
    bullet_panel(
        slide,
        "用户控制",
        [
            "右键菜单设置当前位置。",
            "右键菜单开启 / 关闭 IP 粗定位。",
            "右键菜单开启 / 关闭每日问候。",
            "env.example 保留默认城市与缓存 TTL 配置。",
        ],
        9.15,
        1.65,
        3.35,
        4.85,
    )

    slide = make_slide(prs, "7. 情绪陪伴与唱歌逻辑", "唱歌不再是“难过关键词”的机械触发，而是陪伴之后的选择。")
    flow(
        slide,
        [
            "用户表达情绪",
            "本地分类：轻度 / 明显 / 高风险",
            "轻度：先安慰，不唱歌",
            "明显：先安慰，并邀请是否想听歌",
            "用户回复“唱吧 / 好 / 想听”",
            "播放本地 WAV 歌曲",
        ],
        0.95,
        1.6,
        5.2,
        step_h=0.48,
        gap=0.17,
        color=ACCENT,
    )
    bullet_panel(
        slide,
        "设计原则",
        [
            "把唱歌作为安慰方式之一，而不是固定结果。",
            "高风险表达优先进入谨慎回应，不进入唱歌分支。",
            "同一段低落对话有唱歌提议冷却，避免反复打扰。",
            "Windows 使用 winsound；macOS 走 Qt QSoundEffect 播放 WAV。",
        ],
        6.75,
        1.6,
        5.55,
        4.85,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "8. 技术路线", "整体采用“桌宠 UI + 本地服务 + 外部模型 + 本地数据”的架构。")
    swimlane(
        slide,
        [
            ("桌宠层", "PyQt6：透明窗口、头像、气泡、按钮、右键菜单、沉浸写作窗口。"),
            ("本地服务", "Flask：init、message、weather、memory、creative_doc_stream 等 API。"),
            ("AI 能力", "DeepSeek / OpenAI 兼容接口：对话、创作、总结与视觉扩展。"),
            ("语音链路", "GPT-SoVITS / 系统语音：语音回放、系统听写、可打断播报。"),
            ("数据层", "SQLite：用户偏好、长期摘要、对话回合、每日问候状态。"),
            ("分发层", "Windows exe + start_all；macOS 源码启动包 + start_all.command。"),
        ],
        0.95,
        1.7,
        11.3,
    )

    slide = make_slide(prs, "9. 运行流程图：启动", "双平台启动方式不同，但进入桌宠后的业务逻辑保持一致。")
    bullet_panel(
        slide,
        "Windows 启动",
        [
            "start_all.bat 调用 start_all.ps1。",
            "检测 GPU / CPU 并启动 GPT-SoVITS。",
            "等待本地语音服务端口就绪。",
            "启动 VIVY.exe，初始化 Flask 与 SQLite。",
        ],
        0.75,
        1.55,
        4.15,
        5.05,
    )
    flow(
        slide,
        ["用户双击", "启动脚本", "语音服务", "VIVY.exe", "本地 API", "桌宠可交互"],
        5.3,
        1.55,
        2.7,
        step_h=0.5,
        gap=0.18,
    )
    bullet_panel(
        slide,
        "macOS 启动",
        [
            "start_all.command / start_all.sh 进入 VIVY-src。",
            "自动创建 venv 并安装 requirements。",
            "复制 env.example 为 .env。",
            "使用系统语音与同一套桌宠源码启动。",
        ],
        8.4,
        1.55,
        4.05,
        5.05,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "10. 运行流程图：一次交互", "一次用户输入会先走本地能力分流，再决定是否请求模型。")
    flow(
        slide,
        [
            "用户输入文字 / 语音 / 图片",
            "前端触摸时间与中断旧任务",
            "本地分流：时间、天气、唱歌、停止唱歌",
            "需要 AI 时请求 Flask API",
            "后端读取记忆和运行时上下文",
            "调用模型或流式接口",
            "前端显示、播报、刷新记忆",
        ],
        0.95,
        1.45,
        5.55,
        step_h=0.45,
        gap=0.13,
        color=ACCENT,
    )
    bullet_panel(
        slide,
        "关键控制点",
        [
            "本地功能优先，减少简单问题对模型的依赖。",
            "时间和天气用结构化逻辑处理，稳定性更高。",
            "流式回复支持边生成边显示，并可打断。",
            "语音回放使用 token 机制，避免旧回答继续播放。",
            "每次交互后刷新记忆面板，保持状态可见。",
        ],
        7.0,
        1.45,
        5.25,
        5.15,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "11. 技术方案与数据安全", "项目以本地优先为原则，用户数据默认留在本机。")
    bullet_panel(
        slide,
        "后端接口方案",
        [
            "/api/init：初始化用户会话。",
            "/api/message / message_stream：普通与流式对话。",
            "/api/weather / weather/greeting：天气与问候。",
            "/api/memory：查看、编辑、删除回合。",
            "/api/creative_doc_stream：文档创作辅助。",
        ],
        0.85,
        1.55,
        5.35,
        5.1,
    )
    bullet_panel(
        slide,
        "本地数据方案",
        [
            "vivy.sqlite 保存记忆、偏好和对话回合。",
            ".desktop_user_id 标记本机用户身份。",
            ".env 保存 API Key，分发时不应放入压缩包。",
            "打包给别人时排除 vivy.sqlite 与 .desktop_user_id，保证记忆从零开始。",
            "每日问候状态只保存最近若干天，避免偏好字段无限增长。",
        ],
        6.65,
        1.55,
        5.75,
        5.1,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "12. 双平台交付方案", "当前采用 Windows 成品包 + macOS 源码启动包的阶段性双平台方案。")
    bullet_panel(
        slide,
        "Windows 交付",
        [
            "包含 VIVY.exe / VIVY_layoutupdate.exe。",
            "包含 start_all.bat 和 start_all.ps1 一键启动。",
            "可随包携带 GPT-SoVITS 运行环境。",
            "适合直接压缩分发给普通用户。",
        ],
        0.85,
        1.65,
        5.35,
        4.9,
    )
    bullet_panel(
        slide,
        "macOS 交付",
        [
            "提供 VIVY-src 源码启动包。",
            "start_all.command 自动准备 venv 和依赖。",
            "默认使用系统语音，避免额外语音服务部署复杂度。",
            "后续可继续推进 .app 打包、签名与公证。",
        ],
        6.65,
        1.65,
        5.75,
        4.9,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "13. 创新性", "创新不只是功能堆叠，而是让 AI 以更生活化的方式长期存在。")
    bullet_panel(
        slide,
        "体验创新",
        [
            "把 AI 做成可拖拽、可待机、会主动问候的桌面角色。",
            "通过每日打开次数和时段记忆，减少重复打扰。",
            "情绪陪伴先理解再行动，唱歌变成用户选择，而非关键词触发。",
            "天气、时间、语音和创作能力在一个桌宠入口中自然融合。",
        ],
        0.85,
        1.55,
        5.75,
        5.1,
    )
    bullet_panel(
        slide,
        "技术创新",
        [
            "本地 UI 与 Flask 服务解耦，便于持续扩展功能。",
            "运行时上下文注入，让模型知道当前时间但不污染长期记忆。",
            "结构化偏好 JSON 保存用户状态和每日问候记录。",
            "跨平台音频播放和启动脚本降低部署门槛。",
        ],
        6.9,
        1.55,
        5.45,
        5.1,
        fill=PANEL_2,
    )

    slide = make_slide(prs, "14. 总结与后续方向", "VIVY 已经具备可演示、可分发、可继续迭代的产品雏形。")
    bullet_panel(
        slide,
        "当前成果",
        [
            "完成桌宠 UI、对话、语音、唱歌、记忆、天气、问候和创作辅助闭环。",
            "完成 Windows exe 打包与 macOS 源码启动包同步。",
            "形成本地记忆、每日问候状态和跨会话防重复机制。",
            "具备面向用户演示和小范围分发的基础条件。",
        ],
        0.85,
        1.55,
        5.75,
        5.1,
    )
    bullet_panel(
        slide,
        "后续方向",
        [
            "完善 macOS 原生 .app 打包与签名。",
            "优化天气来源、城市解析和异常网络兜底。",
            "继续扩展更多个性化记忆和长期陪伴策略。",
            "完善发布包清理规则，保护 API Key 和个人记忆隐私。",
        ],
        6.9,
        1.55,
        5.45,
        5.1,
        fill=PANEL_2,
    )

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    print(build())
