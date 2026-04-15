from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(r"d:\lib\CODEGAMES\VIVY\release\VIVY_Flow_Deck.pptx")

BG = RGBColor(8, 15, 25)
PANEL = RGBColor(18, 29, 44)
ACCENT = RGBColor(55, 191, 232)
ACCENT_SOFT = RGBColor(121, 226, 255)
TEXT = RGBColor(240, 248, 255)
TEXT_DIM = RGBColor(180, 201, 220)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.12), Inches(10.5), Inches(0.45))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = TEXT_DIM


def add_bullets(slide, bullets, left=0.95, top=1.7, width=5.1, height=4.8, font_size=20):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = ACCENT_SOFT
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(12)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(9)
    return box


def add_process(slide, steps, left=6.25, top=1.8, width=5.45, step_h=0.55, gap=0.12):
    current_top = top
    for idx, step in enumerate(steps, start=1):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(current_top),
            Inches(width),
            Inches(step_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1.6)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{idx}. {step}"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = TEXT
        current_top += step_h + gap
        if idx < len(steps):
            connector = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + width / 2 - 0.18),
                Inches(current_top - gap + 0.01),
                Inches(0.36),
                Inches(gap + 0.04),
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = ACCENT
            connector.line.color.rgb = ACCENT


def add_architecture(slide):
    labels = [
        ("用户", 0.9),
        ("启动脚本", 2.15),
        ("桌宠程序", 3.4),
        ("本地后端", 4.65),
        ("模型 / 记忆 / 语音", 5.9),
    ]
    left = 3.0
    width = 6.7
    for i, (label, top) in enumerate(labels):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.62),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL if i != 0 else ACCENT
        shape.line.color.rgb = ACCENT_SOFT
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = TEXT
        if i < len(labels) - 1:
            connector = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                Inches(left + width / 2 - 0.22),
                Inches(top + 0.7),
                Inches(0.44),
                Inches(0.42),
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = ACCENT
            connector.line.color.rgb = ACCENT


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(6.9), Inches(11.2), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_DIM


def add_cover(slide):
    set_background(slide)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(0.85),
        Inches(0.18),
        Inches(5.6),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.color.rgb = ACCENT

    add_title(slide, "VIVY 桌宠系统流程图", "启动流程、对话流程、语音流程、记忆流程、创作流程")

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.25),
        Inches(2.0),
        Inches(10.1),
        Inches(3.5),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = ACCENT_SOFT
    box.line.width = Pt(1.8)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    items = [
        "Windows：启动脚本拉起 GPT-SoVITS 与 VIVY.exe",
        "桌宠：接收文字、语音、图片与按钮事件",
        "后端：本地 Flask API 负责对话、记忆与创作接口",
        "AI：调用 DeepSeek / 兼容模型生成回复",
        "数据：本地数据库保存用户记忆与最近回合",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22 if i == 0 else 18)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)
    add_footer(slide, "VIVY 项目流程总览")


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_cover(slide)

    slides = [
        (
            "系统总览",
            "VIVY 整体由启动层、桌宠层、后端层、AI 层、数据层组成。",
            [
                "启动层：start_all.bat / start_all.ps1",
                "语音层：GPT-SoVITS 本地语音服务",
                "桌宠层：VIVY.exe / 双平台桌宠前端",
                "后端层：本地 Flask API",
                "数据层：vivy.sqlite、本地偏好与对话记录",
            ],
            None,
        ),
        (
            "启动流程",
            "从双击脚本到桌宠进入可交互状态。",
            [
                "用户双击 start_all.bat",
                "PowerShell 脚本检测 GPU / CPU",
                "启动 GPT-SoVITS 本地语音服务",
                "等待语音服务端口就绪",
                "启动 VIVY.exe",
                "初始化本地后端、数据库与用户会话",
            ],
            [
                "双击 start_all.bat",
                "调用 start_all.ps1",
                "检测 GPU / CPU",
                "启动 GPT-SoVITS",
                "等待 127.0.0.1:9880 就绪",
                "启动 VIVY.exe",
            ],
        ),
        (
            "对话主流程",
            "用户消息会先经过本地解析，再决定是否走 AI 对话。",
            [
                "输入来源：文字、语音、图片、快捷按钮",
                "先判断是否命中本地唱歌等本地功能",
                "普通消息走 AI 请求",
                "后端读取记忆、偏好、摘要",
                "构造 Prompt 并调用模型",
                "返回文字回复与可选语音播报",
            ],
            [
                "用户输入",
                "桌宠解析请求",
                "请求 Flask API",
                "读取记忆与偏好",
                "调用模型",
                "返回文字 / 语音",
            ],
        ),
        (
            "回复中断机制",
            "这版桌宠支持新问题优先，避免旧回复和新问题串线。",
            [
                "如果上一条回复仍在生成，用户可直接继续提问",
                "新问题发出时，旧语音立即中断",
                "旧流式输出失效，不再刷新界面",
                "新问题接管当前会话",
                "新回复重新开始显示与播报",
            ],
            [
                "旧回复进行中",
                "用户发送新问题",
                "中断旧语音",
                "旧请求失效",
                "新问题接管",
                "新回复显示",
            ],
        ),
        (
            "语音输入流程",
            "语音输入由桌宠弹窗配合系统听写完成。",
            [
                "用户点击“语音”按钮",
                "打开语音输入弹窗",
                "触发系统听写",
                "听写文字落入输入框",
                "用户确认后回填主输入框",
                "进入正常对话流程",
            ],
            [
                "点击语音",
                "打开听写弹窗",
                "系统听写",
                "文字回填",
                "确认发送",
                "进入 AI 对话",
            ],
        ),
        (
            "语音输出流程",
            "回复生成后根据设置决定是否进行语音播报。",
            [
                "检查是否开启语音回放",
                "Windows 优先使用 GPT-SoVITS",
                "macOS 默认使用系统 say",
                "支持流式朗读与完整朗读",
                "新问题到来时可中断当前播报",
            ],
            [
                "回复完成",
                "判断是否播报",
                "进入 TTS",
                "播报进行中",
                "如有新问题则中断",
            ],
        ),
        (
            "本地唱歌流程",
            "唱歌功能属于本地功能，不依赖大模型生成完整歌曲。",
            [
                "识别“唱歌 / 哄我 / 安慰我”等关键词",
                "生成唱歌开场文案",
                "显示气泡并可先播报文案",
                "从 song 目录选择本地 wav",
                "开始播放歌曲",
                "支持“停止唱歌”中断",
            ],
            [
                "识别唱歌关键词",
                "生成开场文案",
                "显示 / 播报提示",
                "选中本地歌曲",
                "开始播放",
                "停止唱歌",
            ],
        ),
        (
            "记忆系统流程",
            "VIVY 会把用户偏好与最近对话写入本地数据库。",
            [
                "记录 user turn 与 assistant turn",
                "保存兴趣信号、形态切换与偏好",
                "定期生成 short summary",
                "定期生成 long summary",
                "后续对话时作为上下文参考",
                "记忆面板支持查看、编辑、删除回合",
            ],
            [
                "用户消息入库",
                "助手回复入库",
                "更新偏好与摘要",
                "写回本地数据库",
                "后续对话复用",
            ],
        ),
        (
            "创作辅助流程",
            "创作模式可读取文档并走流式创作接口。",
            [
                "切换到创作模式",
                "读取文档或打开沉浸写作",
                "调用 creative_doc_stream / office_passage_stream",
                "流式生成润色、续写、点评、加强结果",
                "支持复制、插入和替换文本",
            ],
            [
                "切换创作模式",
                "读取文档 / 选区",
                "构造创作 Prompt",
                "流式返回建议",
                "写作辅助落地",
            ],
        ),
        (
            "图片理解流程",
            "图片会作为结构化消息发送给后端，并按配置决定是否启用视觉模型。",
            [
                "用户选择图片并输入问题",
                "桌宠把图片编码为 data_url",
                "与文字一起发送到 /api/message",
                "后端判断是否配置视觉模型",
                "有视觉模型时进入看图能力",
                "未配置时返回降级提示",
            ],
            [
                "选择图片",
                "编码为 data_url",
                "发送后端",
                "判断视觉模型",
                "看图 / 降级提示",
            ],
        ),
        (
            "双平台说明",
            "Windows 与 macOS 保持相近的交互逻辑，但启动方式不同。",
            [
                "Windows：双击 start_all.bat，直接运行成品包",
                "macOS：执行 start_all.command / start_all.sh，运行源码启动包",
                "Windows 默认本地 GPT-SoVITS 语音服务",
                "macOS 默认系统语音 say",
                "交互逻辑、唱歌文案、状态提示尽量保持一致",
            ],
            [
                "Windows 启动脚本",
                "macOS 启动脚本",
                "统一 UI 与交互",
                "不同 TTS 实现",
                "双平台并行维护",
            ],
        ),
        (
            "总结",
            "VIVY 已形成完整的桌宠 AI 交互闭环。",
            [
                "覆盖启动、对话、语音、记忆、创作与本地唱歌",
                "支持回复中断、状态恢复与双平台兼容",
                "Windows 适合直接发布，macOS 具备源码启动能力",
                "后续仍可继续完善原生打包与 Release 说明",
            ],
            [
                "启动完整",
                "交互完整",
                "记忆完整",
                "创作完整",
                "可继续扩展",
            ],
        ),
    ]

    for title, subtitle, bullets, process in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_title(slide, title, subtitle)
        if title == "系统总览":
            add_bullets(slide, bullets, left=0.8, top=1.8, width=3.9, height=4.9, font_size=17)
            add_architecture(slide)
        else:
            add_bullets(slide, bullets)
            if process:
                add_process(slide, process)
        add_footer(slide, "VIVY 桌宠系统流程图")

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(path)
