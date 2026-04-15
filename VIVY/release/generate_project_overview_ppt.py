from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(r"d:\lib\CODEGAMES\VIVY\release\VIVY_Project_Overview_v2.pptx")

BG = RGBColor(9, 15, 24)
PANEL = RGBColor(18, 28, 43)
PANEL_ALT = RGBColor(24, 38, 58)
ACCENT = RGBColor(56, 190, 232)
ACCENT_SOFT = RGBColor(126, 228, 255)
TEXT = RGBColor(241, 248, 255)
TEXT_DIM = RGBColor(183, 201, 220)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def title_block(slide, title: str, subtitle: str = ""):
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(0.42), Inches(11.8), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.75), Inches(1.02), Inches(11.2), Inches(0.35))
        p2 = sb.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = TEXT_DIM


def footer(slide, text: str = "VIVY 项目汇报"):
    fb = slide.shapes.add_textbox(Inches(0.72), Inches(6.95), Inches(12), Inches(0.22))
    p = fb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = TEXT_DIM


def bullet_panel(slide, title: str, bullets, left, top, width, height, font_size=18):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT_SOFT
    shape.line.width = Pt(1.4)

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(20)
    r0.font.bold = True
    r0.font.color.rgb = ACCENT_SOFT
    p0.space_after = Pt(10)

    for item in bullets:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(7)


def process_stack(slide, title: str, steps, left, top, width):
    head = slide.shapes.add_textbox(Inches(left), Inches(top - 0.45), Inches(width), Inches(0.25))
    p = head.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = ACCENT_SOFT

    current_top = top
    for i, step in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(current_top),
            Inches(width),
            Inches(0.56),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_ALT if i % 2 else PANEL
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.2)
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = f"{i + 1}. {step}"
        r2.font.size = Pt(16)
        r2.font.bold = True
        r2.font.color.rgb = TEXT
        current_top += 0.72
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                Inches(left + width / 2 - 0.18),
                Inches(current_top - 0.16),
                Inches(0.36),
                Inches(0.18),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT
            arr.line.color.rgb = ACCENT


def architecture_bars(slide):
    labels = [
        "用户交互层：桌宠界面、按钮、输入框、语音弹窗",
        "应用逻辑层：请求分流、状态控制、中断机制、唱歌逻辑",
        "本地服务层：Flask API、记忆接口、创作接口",
        "AI 能力层：DeepSeek / 兼容模型、GPT-SoVITS、系统语音",
        "数据层：vivy.sqlite、本地偏好、本地对话记录",
    ]
    top = 1.72
    for idx, label in enumerate(labels):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(2.1),
            Inches(top + idx * 0.95),
            Inches(9.2),
            Inches(0.58),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL if idx % 2 == 0 else PANEL_ALT
        shape.line.color.rgb = ACCENT_SOFT
        shape.line.width = Pt(1.2)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT


def add_slide(prs, title, subtitle, left_title, left_bullets, right_title=None, right_content=None, mode="bullets"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, title, subtitle)
    bullet_panel(slide, left_title, left_bullets, 0.8, 1.7, 5.3, 4.95, font_size=17)
    if right_title and right_content:
        if mode == "process":
            process_stack(slide, right_title, right_content, 6.5, 1.95, 5.7)
        else:
            bullet_panel(slide, right_title, right_content, 6.5, 1.7, 5.7, 4.95, font_size=17)
    footer(slide)
    return slide


def cover(slide):
    set_bg(slide)
    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(0.8),
        Inches(0.18),
        Inches(5.8),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.color.rgb = ACCENT

    title_block(slide, "VIVY 桌宠项目汇报", "作品简介｜产品作用｜用户体验｜功能说明｜技术路线｜运行流程图｜技术方案｜创新性")
    bullet_panel(
        slide,
        "项目定位",
        [
            "VIVY 是一款面向陪伴、对话与创作辅助的 AI 桌面宠物。",
            "项目尝试把桌宠的情绪价值与 AI 的生产力能力结合在一起。",
            "目标是让用户在桌面环境中，以更轻量、更自然的方式使用 AI。",
        ],
        1.2,
        2.0,
        10.2,
        2.4,
        font_size=20,
    )
    bullet_panel(
        slide,
        "汇报重点",
        [
            "作品简介与面向用户的产品价值",
            "功能说明、用户体验与运行流程",
            "技术路线、实现方案与创新性",
        ],
        1.2,
        4.7,
        10.2,
        1.5,
        font_size=18,
    )
    footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cover(slide)

    add_slide(
        prs,
        "作品简介",
        "介绍 VIVY 的定位、目标与主要应用场景。",
        "项目简介",
        [
            "VIVY 是一款桌面 AI 桌宠，目标是把陪伴感、对话感和创作辅助融入桌面常驻形态。",
            "用户无需打开网页，即可通过桌宠完成聊天、获取灵感、调用语音与管理记忆。",
            "项目强调“桌面陪伴 + 本地服务 + AI 能力”的一体化体验。",
        ],
        "应用场景",
        [
            "日常陪伴与情绪交流",
            "写作灵感与创作辅助",
            "语音播报与唱歌互动",
            "桌面常驻式轻交互助手",
        ],
    )

    add_slide(
        prs,
        "产品作用",
        "从用户角度看，VIVY 不只是一个聊天窗口，而是桌面上的陪伴与辅助入口。",
        "产品价值",
        [
            "陪伴作用：通过桌面常驻和轻交互方式，降低用户独处和表达的门槛。",
            "效率作用：把聊天、记忆、语音、创作辅助整合进同一入口，减少应用切换。",
            "情绪作用：支持安慰、唱歌、语音回放等互动方式，增强情绪反馈与陪伴感。",
            "创作作用：通过灵感、续写、润色和文档辅助，为内容创作提供低门槛支持。",
        ],
        "面向用户",
        [
            "希望获得桌面陪伴感的日常用户",
            "需要灵感和写作支持的创作者",
            "希望用更轻量方式与 AI 交流的学生和办公人群",
        ],
    )

    add_slide(
        prs,
        "用户体验",
        "项目在交互设计上强调轻量、连续、可打断和可恢复，尽量贴近真实使用场景。",
        "体验设计重点",
        [
            "低门槛：双击启动即可使用，不需要打开浏览器或复杂配置界面。",
            "连续性：状态提示显示后会自动恢复到主回答，不打断阅读体验。",
            "实时性：用户输入下一个问题时，旧语音和旧回复可立即中断，避免串线。",
            "可控性：支持待机时长、语音回放、记忆模块显示等个性化控制。",
            "陪伴感：语音回放、唱歌、灵感提示等功能让交互更柔和、更有情绪温度。",
        ],
        "用户收益",
        [
            "更自然地使用 AI，而不是被复杂操作打断",
            "在学习、工作、情绪调节中都能快速获得反馈",
            "把桌宠当成日常使用频率更高的个人 AI 入口",
        ],
    )

    add_slide(
        prs,
        "功能说明",
        "VIVY 已形成较完整的桌宠功能闭环。",
        "核心功能",
        [
            "对话：普通聊天、创作模式、流式回复",
            "语音：语音输入、语音回放、可打断播报",
            "本地互动：唱歌、停止唱歌、状态提示",
            "记忆：偏好记录、摘要更新、历史回合管理",
            "创作：文档读取、沉浸写作、Office 选区辅助",
            "视觉：支持图片消息入口与视觉模型扩展",
        ],
        "交互亮点",
        [
            "状态提示显示后可恢复主回答",
            "新问题可中断旧语音与旧回复",
            "待机时长可视化设置",
            "双平台交互逻辑尽量一致",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "技术路线", "项目采用“前端桌宠 + 本地服务 + 外部模型 + 本地数据”的组合路线。")
    architecture_bars(slide)
    bullet_panel(
        slide,
        "技术关键词",
        [
            "PyQt6：桌宠界面与交互",
            "Flask：本地 API 服务",
            "DeepSeek / 兼容模型：对话与创作能力",
            "GPT-SoVITS / 系统 TTS：语音播报",
            "SQLite：本地记忆与对话持久化",
        ],
        0.85,
        5.9,
        12.0,
        0.75,
        font_size=15,
    )
    footer(slide)

    add_slide(
        prs,
        "运行流程图",
        "从启动脚本到桌宠进入可交互状态的主流程。",
        "流程说明",
        [
            "Windows 版先由启动脚本拉起本地语音服务，再启动桌宠程序。",
            "桌宠启动后会初始化本地后端、数据库与用户会话。",
            "初始化完成后，系统进入可交互状态，等待用户输入。",
        ],
        "启动流程",
        [
            "双击 start_all.bat",
            "调用 start_all.ps1",
            "检测 GPU / CPU",
            "启动 GPT-SoVITS",
            "等待端口就绪",
            "启动 VIVY.exe",
            "初始化会话",
        ],
        mode="process",
    )

    add_slide(
        prs,
        "运行流程图",
        "用户发起一次完整交互时，系统的主要处理路径。",
        "交互说明",
        [
            "输入来源可以是文字、语音、图片或快捷按钮。",
            "系统先判断是否命中本地功能，例如唱歌或换个问题。",
            "非本地功能会进入 AI 请求链路，由后端读取记忆和偏好后调用模型。",
            "结果返回后，前端负责显示气泡、语音播报和刷新记忆。",
        ],
        "对话流程",
        [
            "用户输入",
            "判断本地功能",
            "请求本地 API",
            "读取记忆 / 偏好",
            "调用模型",
            "返回回复",
            "显示 / 播报 / 刷新记忆",
        ],
        mode="process",
    )

    add_slide(
        prs,
        "技术方案",
        "VIVY 的实现方案重点在于交互层、服务层与语音链路的协同。",
        "前后端方案",
        [
            "前端采用 PyQt6，负责桌宠形象、输入区、气泡区、记忆面板和创作入口。",
            "后端采用 Flask，提供 init、message、message_stream、memory、creative_doc_stream 等接口。",
            "桌宠与后端通过本地 HTTP 调用，保持 UI 与业务逻辑分离。",
        ],
        "语音与数据方案",
        [
            "Windows 默认本地 GPT-SoVITS，macOS 默认系统 say。",
            "语音输入走系统听写，降低跨平台接入难度。",
            "用户记忆、本地偏好与对话回合写入 SQLite。",
            "通过中断 token 机制保证新问题可打断旧回复。",
        ],
    )

    add_slide(
        prs,
        "技术方案",
        "双平台兼容采取“Windows 成品包 + macOS 源码启动包”的阶段性方案。",
        "Windows 方案",
        [
            "保留完整的本地 GPT-SoVITS 运行环境",
            "通过 start_all.bat / start_all.ps1 一键启动",
            "适合直接分发与开箱即用",
        ],
        "macOS 方案",
        [
            "通过 start_all.command / start_all.sh 启动",
            "自动创建 venv 并安装依赖",
            "使用跨平台桌宠源码和系统语音",
            "在交互逻辑上尽量与 Windows 保持一致",
        ],
    )

    add_slide(
        prs,
        "创新性",
        "项目的创新不只在技术组合，更在于把 AI 能力做成贴近用户桌面的长期陪伴形态。",
        "创新点",
        [
            "将桌宠、对话、记忆、语音、创作辅助融合为统一交互入口，而不是分散在多个工具中。",
            "把“陪伴型产品体验”和“生产力型 AI 能力”放在同一产品中同时实现。",
            "通过回复中断机制、状态恢复机制等细节，让桌宠体验更接近真实对话对象。",
            "采用双平台兼容思路，兼顾 Windows 成品发布与 macOS 启动方案。",
        ],
        "面向用户的意义",
        [
            "让 AI 更像一个长期陪伴的桌面角色，而不只是一次性问答工具。",
            "让用户在低压力、低切换成本的环境里获得帮助。",
            "在功能丰富的同时，尽量保留桌宠交互的轻巧和情绪价值。",
        ],
    )

    add_slide(
        prs,
        "总结",
        "VIVY 已经具备可演示、可交互、可继续扩展的完整桌宠形态。",
        "项目成果",
        [
            "完成桌宠 UI、对话、语音、唱歌、记忆、创作辅助等核心闭环。",
            "已支持双平台运行方案，并完成主要界面逻辑对齐。",
            "形成了从启动、交互到本地数据管理的完整技术链路。",
        ],
        "后续方向",
        [
            "进一步完善 macOS 原生打包",
            "补齐 Release 说明与交付文档",
            "提升流式失败场景的兜底策略",
            "继续优化视觉与语音体验",
        ],
    )

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    print(build())
